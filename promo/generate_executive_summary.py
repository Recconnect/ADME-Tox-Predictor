from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

CHARTS = os.path.join(os.path.dirname(__file__), "charts")
OUTPUT = os.path.join(os.path.dirname(__file__), "ADMET_Predictor_Executive_Summary.pptx")

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_BLUE = RGBColor(0x34, 0x98, 0xDB)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
GRAY_LIGHT = RGBColor(0xEC, 0xF0, 0xF1)
GRAY_MED = RGBColor(0x7F, 0x8C, 0x8D)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_bullet_text(slide, left, top, width, height, items, size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return txBox

def slide_title_right(slide, title, subtitle=None):
    add_bg(slide, DARK)
    add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT)
    add_text(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1), title, size=40, bold=True)
    if subtitle:
        add_text(slide, Inches(1.5), Inches(3), Inches(10), Inches(0.8), subtitle, size=20, color=GRAY_MED)

# ====== SLIDE 1: Title ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.15), ACCENT)
add_shape(slide, Inches(0), H - Inches(0.15), W, Inches(0.15), ACCENT)
add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
         "ADME/Tox Predictor", size=54, bold=True)
add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(1),
         "AI-скрининг ADME-свойств и токсичности на ранних стадиях", size=28, color=GRAY_MED)
add_text(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.6),
         "Презентация для инвесторов  |  Июль 2026", size=20, color=ACCENT_BLUE)
add_shape(slide, Inches(1), Inches(5.0), Inches(3), Inches(0.06), ACCENT)
add_bullet_text(slide, Inches(1), Inches(5.3), Inches(11), Inches(1.5), [
    "MVP готов  •  3 ADME/Tox свойства  •  Метрики превышают цели",
    "Streamlit UI  •  Валидация на известных препаратах  •  Данные Harvard Dataverse"
], size=18, color=GRAY_MED)

# ====== SLIDE 2: Проблема и решение ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_title_right(slide, "Проблема", "90% кандидатов проваливаются в клинике из-за ADME/Tox")
add_bullet_text(slide, Inches(0.6), Inches(1.5), Inches(5.5), Inches(5), [
    "Стоимость разработки препарата: $2.23 млрд",
    "60% бюджета уходит на lead optimization",
    "Eroom's Law: эффективность R&D падает вдвое каждые 9 лет",
    "Скрининг требует синтеза тысяч молекул",
], size=17, color=GRAY_LIGHT)
if os.path.exists(os.path.join(CHARTS, "eroms_law.png")):
    slide.shapes.add_picture(os.path.join(CHARTS, "eroms_law.png"),
                             Inches(6.5), Inches(1.5), Inches(6), Inches(3.5))
add_shape(slide, Inches(6.5), Inches(5.3), Inches(6), Inches(0.06), ACCENT)
add_text(slide, Inches(6.5), Inches(5.5), Inches(6), Inches(1),
         "Решение: AI-модель предсказывает свойства до синтеза — <1 сек на молекулу",
         size=18, color=ACCENT, bold=True)

# ====== SLIDE 3: Метрики ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_title_right(slide, "Метрики", "Три модели LightGBM — все превышают целевые показатели")
if os.path.exists(os.path.join(CHARTS, "metrics_comparison.png")):
    slide.shapes.add_picture(os.path.join(CHARTS, "metrics_comparison.png"),
                             Inches(0.5), Inches(1.5), Inches(8.5), Inches(3.5))

y_start = Inches(5.3)
add_shape(slide, Inches(0.5), y_start, Inches(12.3), Inches(1.8), RGBColor(0x22, 0x22, 0x3A))
headers = ["Свойство", "Датасет", "Размер", "Модель", "Цель", "Достигнуто", "Статус"]
widths = [Inches(2.5), Inches(2), Inches(1.2), Inches(2.2), Inches(1.5), Inches(1.5), Inches(1.5)]
x_pos = Inches(0.7)
for header, w in zip(headers, widths):
    add_text(slide, x_pos, y_start + Inches(0.1), w, Inches(0.5), header, size=13, color=ACCENT_BLUE, bold=True)
    x_pos += w

rows = [
    ("Растворимость (logS)", "AqSolDB", "9 982", "LightGBM регр.", "R² > 0.6", "R² = 0.826", "Превышена"),
    ("Caco-2 проницаемость", "Wang et al.", "910", "LightGBM клас.", "Acc > 75%", "Acc = 83.0%", "Превышена"),
    ("hERG токсичность", "TDC", "655", "LightGBM клас.", "AUC > 0.7", "AUC = 0.873", "Превышена"),
]
for r_idx, row in enumerate(rows):
    y = y_start + Inches(0.65 + r_idx * 0.4)
    x_pos = Inches(0.7)
    for val, w in zip(row, widths):
        c = ACCENT if val == "Превышена" else GRAY_LIGHT
        b = True if val == "Превышена" else False
        add_text(slide, x_pos, y, w, Inches(0.4), val, size=12, color=c, bold=b)
        x_pos += w

# ====== SLIDE 4: Рынок ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_title_right(slide, "Рынок", "AI в drug discovery: $3.5 млрд (2025), CAGR 30%")
if os.path.exists(os.path.join(CHARTS, "market_growth.png")):
    slide.shapes.add_picture(os.path.join(CHARTS, "market_growth.png"),
                             Inches(0.5), Inches(1.5), Inches(7), Inches(4))
items = [
    "TAM: $3.5 млрд (2025) → $7.6 млрд (2028)",
    "CAGR: 30% — фарма активно внедряет AI",
    "Ниша: ADME/Tox — слабо занята",
    "Клиенты: CRO, биотех, R&D отделы фармы",
    "Модель: SaaS/API $10K–50K в год за лицензию",
]
add_bullet_text(slide, Inches(8), Inches(1.5), Inches(4.8), Inches(5), items, size=17, color=GRAY_LIGHT)

# ====== SLIDE 5: Дорожная карта ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_title_right(slide, "Дорожная карта", "От MVP до Enterprise SaaS")
phases = [
    ("MVP (Готово)", "2-4 недели", [
        "3 модели LightGBM",
        "Streamlit веб-интерфейс",
        "Валидация на 10 препаратах",
    ], ACCENT),
    ("Пилот", "2-3 месяца", [
        "GNN (+5-10% к метрикам)",
        "API для интеграции",
        "Пилот с 1-2 CRO",
    ], ACCENT_BLUE),
    ("Масштаб", "6-12 месяцев", [
        "5+ ADME/Tox свойств",
        "Enterprise SaaS платформа",
        "Собственный пайплайн данных",
    ], ACCENT_RED),
]
for i, (title, timeline, items, color) in enumerate(phases):
    x = Inches(1 + i * 4.2)
    y = Inches(2.0)
    add_shape(slide, x, y, Inches(3.6), Inches(0.06), color)
    add_text(slide, x, y + Inches(0.2), Inches(3.6), Inches(0.5), title, size=22, bold=True, color=color)
    add_text(slide, x, y + Inches(0.7), Inches(3.6), Inches(0.4), timeline, size=14, color=GRAY_MED)
    add_bullet_text(slide, x + Inches(0.1), y + Inches(1.2), Inches(3.4), Inches(3.5),
                    [f"  {it}" for it in items], size=14, color=GRAY_LIGHT)

# ====== SLIDE 6: Запрос и контакты ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.15), ACCENT)
add_shape(slide, Inches(0), H - Inches(0.15), W, Inches(0.15), ACCENT)
add_text(slide, Inches(1), Inches(1.0), Inches(11), Inches(1),
         "Запрос: Seed $200K–$500K", size=40, bold=True, align=PP_ALIGN.CENTER)
items_left = [
    "Fine-tuning на корпоративных данных: $100K",
    "Наём 2-3 специалистов: $150K",
    "Облачная инфраструктура + GPU: $50K",
    "Маркетинг и продажи: $50–100K",
]
items_right = [
    "Цель: пилот с 2-3 фарма/CRO партнёрами",
    "Закрытие раунда: Q4 2026",
    "Следующий раунд: Series A $5–10M",
]
add_bullet_text(slide, Inches(1.5), Inches(2.5), Inches(4.5), Inches(4), items_left, size=18, color=GRAY_LIGHT)
add_bullet_text(slide, Inches(7), Inches(2.5), Inches(4.5), Inches(4), items_right, size=18, color=GRAY_LIGHT)
add_shape(slide, Inches(6.5), Inches(2.5), Inches(0.06), Inches(4), ACCENT)

add_shape(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.06), GRAY_MED)
add_text(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
         "ADME/Tox Predictor  |  MVP готов  |  Открыты к пилотным проектам  |  D:\\AI\\biotech\\adme_proto",
         size=16, color=GRAY_MED, align=PP_ALIGN.CENTER)

prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
