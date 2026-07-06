from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

CHARTS = os.path.join(os.path.dirname(__file__), "charts")
OUTPUT = os.path.join(os.path.dirname(__file__), "ADMET_Predictor_Investor_Pitch.pptx")

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_BLUE = RGBColor(0x34, 0x98, 0xDB)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
GRAY_LIGHT = RGBColor(0xEC, 0xF0, 0xF1)
GRAY_MED = RGBColor(0x7F, 0x8C, 0x8D)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_bullet_text(slide, left, top, width, height, items, size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return txBox

def add_subtitle_line(slide, left, top, width=Inches(0.06), height=Inches(0.5), color=ACCENT):
    return add_shape(slide, left, top, width, height, color)

def slide_title_right(slide, title, subtitle=None):
    add_bg(slide, DARK)
    add_shape(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT)
    add_text(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1), title, size=40, bold=True)
    if subtitle:
        add_text(slide, Inches(1.5), Inches(3), Inches(10), Inches(0.8), subtitle, size=20, color=GRAY_MED)

# ====== SLIDE 1: Title ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.15), ACCENT)
add_shape(slide, Inches(0), H - Inches(0.15), W, Inches(0.15), ACCENT)
add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
         "ADME/Tox Predictor", size=54, bold=True)
add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(1),
         "AI-powered Early-Stage Drug ADME/Tox Screening", size=28, color=GRAY_MED)
add_text(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.6),
         "Investor Pitch  |  July 2026", size=20, color=ACCENT_BLUE)
add_shape(slide, Inches(1), Inches(5.0), Inches(3), Inches(0.06), ACCENT)
add_bullet_text(slide, Inches(1), Inches(5.3), Inches(11), Inches(1.5), [
    "MVP ready  •  3 ADME/Tox properties  •  Exceeded target metrics",
    "Streamlit UI  •  Validation on known drugs  •  Harvard Dataverse data"
], size=18, color=GRAY_MED)

# ====== SLIDE 2: The Problem ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_title_right(slide, "The Problem", "Why 90% of drug candidates fail in clinical trials")
if os.path.exists(os.path.join(CHARTS, "eroms_law.png")):
    slide.shapes.add_picture(os.path.join(CHARTS, "eroms_law.png"),
                             Inches(0.6), Inches(1.2), Inches(7), Inches(4))
add_bullet_text(slide, Inches(8.2), Inches(1.2), Inches(4.5), Inches(5.5), [
    "90% fail rate — primary cause: poor ADME/Tox",
    "$2.23B average drug development cost",
    "Eroom's Law: R&D efficiency halves every 9 years",
    "60% of budget wasted on wrong molecules",
    "Traditional screening: slow, expensive, low throughput"
], size=17, color=GRAY_LIGHT)

# ====== SLIDE 3: The Solution ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_title_right(slide, "The Solution", "AI-powered ADME/Tox prediction before synthesis")
if os.path.exists(os.path.join(CHARTS, "business_case.png")):
    slide.shapes.add_picture(os.path.join(CHARTS, "business_case.png"),
                             Inches(0.6), Inches(1.5), Inches(7.5), Inches(3.8))
items = [
    "<1 sec prediction per molecule on CPU",
    "54mo -> 18mo discovery timeline",
    "$2.23B -> ~$220M development cost",
    "5-10x fewer molecules to synthesize",
    "Validated on known drugs (Aspirin, Ibuprofen, etc.)"
]
add_bullet_text(slide, Inches(8.5), Inches(1.5), Inches(4.2), Inches(5), items, size=17, color=GRAY_LIGHT)

# ====== SLIDE 4: Metrics ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_title_right(slide, "Model Performance", "MVP metrics across 3 ADME/Tox properties")
if os.path.exists(os.path.join(CHARTS, "metrics_comparison.png")):
    slide.shapes.add_picture(os.path.join(CHARTS, "metrics_comparison.png"),
                             Inches(0.5), Inches(1.5), Inches(8.5), Inches(3.5))

# Metrics table as shapes
y_start = Inches(5.3)
add_shape(slide, Inches(0.5), y_start, Inches(12.3), Inches(1.8), RGBColor(0x22, 0x22, 0x3A))
headers = ["Property", "Dataset", "Size", "Model", "Target", "Achieved", "Status"]
widths = [Inches(2.5), Inches(2), Inches(1.2), Inches(2.2), Inches(1.5), Inches(1.5), Inches(1.5)]
x_pos = Inches(0.7)
for header, w in zip(headers, widths):
    add_text(slide, x_pos, y_start + Inches(0.1), w, Inches(0.5), header, size=13, color=ACCENT_BLUE, bold=True)
    x_pos += w

rows = [
    ("Solubility (logS)", "AqSolDB", "9,982", "LightGBM reg.", "R² > 0.6", "R² = 0.826", "Exceeded"),
    ("Caco-2 Permeability", "Wang et al.", "910", "LightGBM cls.", "Acc > 75%", "Acc = 83.0%", "Exceeded"),
    ("hERG Toxicity", "TDC", "655", "LightGBM cls.", "AUC > 0.7", "AUC = 0.873", "Exceeded"),
]
for r_idx, row in enumerate(rows):
    y = y_start + Inches(0.65 + r_idx * 0.4)
    x_pos = Inches(0.7)
    for val, w in zip(row, widths):
        c = ACCENT if val == "Exceeded" else GRAY_LIGHT
        b = True if val == "Exceeded" else False
        add_text(slide, x_pos, y, w, Inches(0.4), val, size=12, color=c, bold=b)
        x_pos += w

# ====== SLIDE 5: Market ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_title_right(slide, "Market Opportunity", "AI in Drug Discovery: $3.5B in 2025, growing at 30% CAGR")
if os.path.exists(os.path.join(CHARTS, "market_growth.png")):
    slide.shapes.add_picture(os.path.join(CHARTS, "market_growth.png"),
                             Inches(0.5), Inches(1.5), Inches(7), Inches(4))
items = [
    "TAM: $3.5B (2025) -> $7.6B (2028)",
    "CAGR: 30% driven by pharma AI adoption",
    "Target: ADME/Tox niche — underserved",
    "Clients: CROs, biotech, pharma R&D",
    "Business model: SaaS/API $10K-50K/yr"
]
add_bullet_text(slide, Inches(8), Inches(1.5), Inches(4.8), Inches(5), items, size=17, color=GRAY_LIGHT)

# ====== SLIDE 6: Lead Optimization ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_title_right(slide, "Why ADME/Tox?", "60% of R&D budget goes to Lead Optimization — our target")
if os.path.exists(os.path.join(CHARTS, "lead_optimization.png")):
    slide.shapes.add_picture(os.path.join(CHARTS, "lead_optimization.png"),
                             Inches(0.5), Inches(1.5), Inches(6), Inches(4.5))
items = [
    "60% of pharma R&D = lead optimization",
    "Molecules with poor ADME/Tox screened out",
    "Our model: 5-10x fewer syntheses needed",
    "ROI for client: 10x-50x on subscription",
    "Savings: $1-5M per drug program"
]
add_bullet_text(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(5), items, size=17, color=GRAY_LIGHT)

# ====== SLIDE 7: Roadmap ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_title_right(slide, "Roadmap", "From MVP to Enterprise SaaS")
phases = [
    ("MVP (Done)", "2-4 weeks", [
        "3 LightGBM models (solubility, Caco-2, hERG)",
        "Streamlit web UI",
        "Validation on 10 known drugs",
    ], ACCENT),
    ("Pilot", "2-3 months", [
        "GNN +5-10% improvement",
        "API development",
        "Pilot with 1-2 CRO partners",
    ], ACCENT_BLUE),
    ("Scale", "6-12 months", [
        "5+ ADME/Tox properties",
        "Enterprise SaaS platform",
        "Proprietary data pipeline",
    ], ACCENT_RED),
]
for i, (title, timeline, items, color) in enumerate(phases):
    x = Inches(1 + i * 4.2)
    y = Inches(2.0)
    add_shape(slide, x, y, Inches(3.6), Inches(0.06), color)
    add_text(slide, x, y + Inches(0.2), Inches(3.6), Inches(0.5), title, size=22, bold=True, color=color)
    add_text(slide, x, y + Inches(0.7), Inches(3.6), Inches(0.4), timeline, size=14, color=GRAY_MED)
    add_bullet_text(slide, x + Inches(0.1), y + Inches(1.2), Inches(3.4), Inches(3.5),
                    [f"  {it}" for it in items], size=14, color=GRAY_LIGHT)

# ====== SLIDE 8: Ask ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_title_right(slide, "Funding Ask", "Seed round: $200K - $500K")
items_left = [
    "Fine-tuning on corporate data: $100K",
    "Hiring (2-3 specialists): $150K",
    "Cloud GPU + data infra: $50K",
    "Marketing & sales: $50-100K",
]
items_right = [
    "Goal: pilot with 2-3 pharma/CRO partners",
    "Target close: Q4 2026",
    "Next round: Series A at $5-10M valuation",
]
add_bullet_text(slide, Inches(1.5), Inches(2.5), Inches(4.5), Inches(4), items_left, size=18, color=GRAY_LIGHT)
add_bullet_text(slide, Inches(7), Inches(2.5), Inches(4.5), Inches(4), items_right, size=18, color=GRAY_LIGHT)
add_shape(slide, Inches(6.5), Inches(2.5), Inches(0.06), Inches(4), ACCENT)

# ====== SLIDE 9: Why Now ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_title_right(slide, "Why Now?", "The perfect storm for AI in drug discovery")
items = [
    "Market growing: $3.5B -> $7.6B in 3 years",
    "Data is available: TDC, ChEMBL, PubChem",
    "Insilico proved the model: Rentosertib in Phase IIa",
    "Eroom's Law makes AI a necessity, not a luxury",
    "MVP is already running with exceeded metrics",
    "Low capital requirement to reach pilot stage",
]
add_bullet_text(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(4), items, size=20, color=GRAY_LIGHT)

# ====== SLIDE 10: Contact ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_shape(slide, Inches(0), Inches(H/2 - 0.075), W, Inches(0.15), ACCENT)
add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
         "ADME/Tox Predictor", size=48, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
         "Thank you", size=32, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.8),
         "MVP ready | Open to pilot projects | Looking for partners",
         size=18, color=GRAY_MED, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.8),
         "D:\\AI\\biotech\\adme_proto",
         size=16, color=GRAY_MED, align=PP_ALIGN.CENTER)

prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
